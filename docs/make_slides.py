"""Generate presentation slides for ISBA 4715 capstone."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DARK = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)
GRAY = RGBColor(0x5D, 0x6D, 0x7E)
GREEN = RGBColor(0x27, 0xAE, 0x60)
WIDTH = Inches(13.333)
HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return tf


def make_slides():
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)

    add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 "MLB Player Performance Analytics",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
                 "Where Traditional Stats and Statcast Metrics Agree, Disagree,\nand What It Means for Player Evaluation",
                 font_size=22, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
                 "AJ Schoolcraft  |  ISBA 4715 Capstone  |  Spring 2026",
                 font_size=18, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.5),
                 "Target Role: Baseball Operations Analyst, San Francisco Giants",
                 font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

    # --- Slide 2: Problem & Analytical Question ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Can Statcast Data Identify Undervalued MLB Hitters?",
                 font_size=32, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "The Problem with Traditional Batting Stats",
                 font_size=16, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), [
        "Batting average treats all hits equally and ignores walks,\nHBP, and quality of contact",
        "A .250 hitter crushing the ball at 100+ mph looks the\nsame as a .250 hitter blooping singles",
        "Since 2015, Statcast tracks exit velocity, launch angle,\nand barrel rate — measuring HOW players hit, not just IF",
        "xwOBA uses this data to estimate what a player's stats\nSHOULD be based on contact quality",
        "The gap between xwOBA and traditional stats reveals\nwhich players are undervalued or overvalued",
    ], font_size=16)

    # Data scope box
    shape = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Data Scope"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER

    for line in [
        "",
        "2024–2025 MLB seasons",
        "All 30 teams, ~800 qualified batters",
        "",
        "Sources:",
        "  MLB Stats API (bio, game logs, season stats)",
        "  Pybaseball / Statcast (exit velo, barrel %, xwOBA)",
        "",
        "16 research sources scraped from FanGraphs,",
        "Baseball Savant, Baseball Prospectus, Driveline",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

    # --- Slide 3: Descriptive Insight ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Traditional Stats Systematically Undervalue Elite Contact Quality",
                 font_size=32, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "Descriptive: What happened across the league?",
                 font_size=16, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), [
        "Players with barrel rates above 10% outperform their batting\naverage by ~45 points in xwOBA",
        "Barrels produce a .684 BA and 2.236 SLG vs. .292/.368\nfor non-barreled contact (2026 data)",
        "The scatter plot reveals a clear cluster of 'undervalued'\nhitters — high xwOBA, low batting average",
        "These gaps are not random — they correlate with\nmeasurable batted-ball quality (exit velocity, barrel %)",
    ], font_size=16)

    # Placeholder for screenshot
    shape = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "\n\n\n[Screenshot: League Overview scatter plot\nfrom Tab 1 of dashboard]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # --- Slide 3: Diagnostic Insight ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "High Exit Velocity + Low BABIP = Unlucky, Not Untalented",
                 font_size=32, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "Diagnostic: Why do results diverge from underlying quality?",
                 font_size=16, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), [
        "The gap between xwOBA and batting average is largest for\nplayers with high exit velocity but below-average BABIP",
        "Rolling 15-game trends reveal when the gap develops —\noften after a streak of hard-hit balls finding fielders",
        "Contact quality charts (barrel % + exit velo) confirm the\nhitter's process is sound even when results lag",
        "BABIP regresses toward ~.300 over time — players with\nlow BABIP and high barrel rates are buy-low candidates",
    ], font_size=16)

    shape = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "\n\n\n[Screenshot: Player Deep Dive trend chart\nfrom Tab 2 of dashboard]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # --- Slide 4: Recommendation ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "Recommendation: Target Undervalued Hitters for Acquisition",
                 font_size=32, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
                 "Action  →  Expected Outcome",
                 font_size=18, color=ACCENT, bold=True)

    add_bullet_list(slide, Inches(1.2), Inches(2.2), Inches(10.5), Inches(1.5), [
        "ACTION: Target players with top-quartile barrel rates (>10%) and below-average BABIP (<.280)",
        "OUTCOME: Expected regression to the mean projects 15–25 point batting average increases",
        "VALUE: These players are available below market value because traditional stats understate their talent",
    ], font_size=18)

    # Decision framework box
    shape = slide.shapes.add_shape(
        1, Inches(1.2), Inches(4.2), Inches(10.5), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Decision Framework"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK

    for line in [
        "1. Screen: Barrel % > 10%, BABIP < .280, PA > 100",
        "2. Diagnose: Check xwOBA vs. BA gap score (positive = undervalued)",
        "3. Confirm: Review rolling trend — is contact quality sustained or fluky?",
        "4. Act: Acquire before BABIP regression corrects the market price",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_before = Pt(6)

    # --- Slide 5: Architecture ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 "End-to-End Pipeline: From Raw APIs to Actionable Dashboard",
                 font_size=32, color=DARK, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "Architecture & Tech Stack",
                 font_size=16, color=ACCENT, bold=True)

    # Left column - pipeline description
    add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0), [
        "Extract: MLB Stats API + Pybaseball (Statcast)\npull player stats, game logs, and batted-ball data",
        "Load: Python scripts write to Snowflake RAW schema\nvia snowflake-connector-python",
        "Transform: dbt builds staging views and mart tables\n(star schema: 2 facts + 4 dimensions)",
        "Present: Streamlit dashboard with Plotly charts\ndeployed to Community Cloud",
        "Knowledge Base: 16 sources scraped from FanGraphs,\nBaseball Savant, BP → 4 synthesized wiki pages",
        "Orchestrate: GitHub Actions runs daily ELT\nand weekly knowledge base re-scrape",
    ], font_size=15)

    shape = slide.shapes.add_shape(
        1, Inches(7.0), Inches(1.9), Inches(5.5), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = "\n\n\n[Screenshot: Pipeline diagram\nfrom README on GitHub]"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # --- Slide 7: Conclusion & Next Steps ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                 "Key Takeaways",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(2.5), [
        "Traditional stats systematically undervalue hitters with elite contact quality",
        "The xwOBA vs. batting average gap is a reliable signal for identifying\nundervalued players — driven by exit velocity, barrel rate, and BABIP variance",
        "A data-driven acquisition framework using Statcast metrics can surface\nabove-market-value talent before traditional stats catch up",
    ], font_size=18, color=WHITE)

    add_text_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
                 "Future Enhancements",
                 font_size=22, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(2.0), [
        "Add pitching analysis tab (FIP vs. ERA gaps for pitcher evaluation)",
        "Normalize gap score to z-scores for cross-season comparison",
        "Integrate sprint speed and defensive metrics (OAA) for full-player valuation",
    ], font_size=16, color=RGBColor(0xAA, 0xBB, 0xCC))

    # --- Save ---
    out = "docs/slides.pptx"
    prs.save(out)
    print(f"Saved to {out}")
    print("Next: Open in PowerPoint/Google Slides, add screenshots, export as PDF")


if __name__ == "__main__":
    make_slides()
